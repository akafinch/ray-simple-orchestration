#!/usr/bin/env python3
"""Single architecture slide — light mode, Akamai brand. Speaker notes, not on-slide bullets."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Akamai brand
AKAMAI_BLUE = RGBColor(0x00, 0x99, 0xCC)
AKAMAI_ORANGE = RGBColor(0xFF, 0x99, 0x33)
NAVY = RGBColor(0x00, 0x28, 0x56)
CHARCOAL = RGBColor(0x54, 0x56, 0x5B)
GRAY = RGBColor(0xA8, 0xA8, 0xAA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF7, 0xF8, 0xFA)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BLUE = RGBColor(0xE8, 0xF4, 0xFA)
CARD_ORANGE = RGBColor(0xFE, 0xF3, 0xE5)
CARD_GRAY = RGBColor(0xF0, 0xF1, 0xF3)
BORDER_LIGHT = RGBColor(0xDE, 0xE0, 0xE3)

FONT = "Roboto"
FONT_MONO = "Consolas"


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def tb(slide, l, t, w, h, text, sz=14, color=CHARCOAL, bold=False,
       align=PP_ALIGN.LEFT, font=FONT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tf


def card(slide, l, t, w, h, fill=CARD_BG, border=BORDER_LIGHT):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(1)
    s.shadow.inherit = False
    return s


def accent_bar(slide, l, t, w, color=AKAMAI_BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(3))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

# ==========================================
# Header
# ==========================================
accent_bar(slide, Inches(0), Inches(0), Inches(13.333), AKAMAI_BLUE)
tb(slide, Inches(0.6), Inches(0.2), Inches(5), Inches(0.4),
   "AKAMAI CLOUD COMPUTING", sz=12, color=AKAMAI_BLUE, bold=True)
tb(slide, Inches(0.6), Inches(0.55), Inches(10), Inches(0.6),
   "CLIP/CLAP GPU Inference \u2014 System Architecture", sz=26, color=NAVY, bold=True)
accent_bar(slide, Inches(0.6), Inches(1.1), Inches(2.5))

# ==========================================
# LKE-E cluster boundary
# ==========================================
card(slide, Inches(0.4), Inches(1.4), Inches(12.5), Inches(4.5), CARD_GRAY, BORDER_LIGHT)
tb(slide, Inches(0.7), Inches(1.5), Inches(8), Inches(0.35),
   "LKE Enterprise Cluster  \u2014  Private VPC", sz=12, color=CHARCOAL, bold=True, font=FONT_MONO)

# --- CPU Pool ---
card(slide, Inches(0.7), Inches(1.95), Inches(4.8), Inches(3.7), CARD_BLUE, AKAMAI_BLUE)
tb(slide, Inches(0.95), Inches(2.05), Inches(4), Inches(0.3),
   "CPU NODE POOL", sz=12, color=AKAMAI_BLUE, bold=True, font=FONT_MONO)
tb(slide, Inches(0.95), Inches(2.3), Inches(4), Inches(0.25),
   "2x g6-standard-4", sz=10, color=GRAY, font=FONT_MONO)

cpu_items = [
    ("Ray Head", "Scheduler"),
    ("KubeRay Operator", "Lifecycle mgmt"),
    ("Demo App", "FastAPI + UI"),
    ("Prometheus", "Metrics"),
    ("Grafana", "Dashboards"),
]
for i, (name, desc) in enumerate(cpu_items):
    col, row = i % 3, i // 3
    x = Inches(0.85 + col * 1.55)
    y = Inches(2.7 + row * 1.4)
    card(slide, x, y, Inches(1.4), Inches(1.15), CARD_BG, BORDER_LIGHT)
    tb(slide, x + Inches(0.1), y + Inches(0.1), Inches(1.2), Inches(0.3),
       name, sz=11, color=NAVY, bold=True, font=FONT_MONO)
    tb(slide, x + Inches(0.1), y + Inches(0.45), Inches(1.2), Inches(0.5),
       desc, sz=10, color=CHARCOAL)

# --- GPU Pool ---
card(slide, Inches(5.8), Inches(1.95), Inches(6.8), Inches(3.7), CARD_ORANGE, AKAMAI_ORANGE)
tb(slide, Inches(6.05), Inches(2.05), Inches(6), Inches(0.3),
   "GPU NODE POOL  \u2014  3x RTX 4000 Ada", sz=12, color=AKAMAI_ORANGE, bold=True, font=FONT_MONO)
tb(slide, Inches(6.05), Inches(2.3), Inches(6), Inches(0.25),
   "g2-gpu-rtx4000a1-s  |  NoSchedule taint", sz=10, color=GRAY, font=FONT_MONO)

gpu_workers = [
    ("Worker 1", "CLIP Replica 1", "768-dim text+image"),
    ("Worker 2", "CLIP Replica 2", "768-dim text+image"),
    ("Worker 3", "CLAP Replica 1", "512-dim audio"),
]
for i, (name, model, dims) in enumerate(gpu_workers):
    x = Inches(6.0 + i * 2.15)
    card(slide, x, Inches(2.7), Inches(2.0), Inches(2.0), CARD_BG, BORDER_LIGHT)
    tb(slide, x + Inches(0.15), Inches(2.8), Inches(1.7), Inches(0.3),
       name, sz=12, color=NAVY, bold=True, font=FONT_MONO)
    tb(slide, x + Inches(0.15), Inches(3.1), Inches(1.7), Inches(0.3),
       model, sz=10, color=CHARCOAL)
    tb(slide, x + Inches(0.15), Inches(3.4), Inches(1.7), Inches(0.3),
       dims, sz=9, color=GRAY, font=FONT_MONO)
    # GPU badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x + Inches(0.4), Inches(3.85), Inches(1.2), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = AKAMAI_ORANGE
    badge.line.fill.background()
    tb(slide, x + Inches(0.4), Inches(3.86), Inches(1.2), Inches(0.32),
       "1x GPU", sz=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_MONO)

# Init container note
card(slide, Inches(6.0), Inches(4.85), Inches(6.5), Inches(0.5), CARD_BG, BORDER_LIGHT)
tb(slide, Inches(6.2), Inches(4.88), Inches(6.1), Inches(0.45),
   "Init containers pull model weights from Object Storage at pod startup",
   sz=10, color=CHARCOAL)

# ==========================================
# External components — bottom row
# ==========================================

# Cloud Firewall + NodeBalancer
card(slide, Inches(0.4), Inches(6.15), Inches(5.2), Inches(0.45), AKAMAI_BLUE, AKAMAI_BLUE)
tb(slide, Inches(0.6), Inches(6.17), Inches(4.8), Inches(0.4),
   "Cloud Firewall  \u2014  IP-restricted access", sz=11, color=WHITE, bold=True,
   align=PP_ALIGN.CENTER)

card(slide, Inches(0.4), Inches(6.65), Inches(5.2), Inches(0.55), CARD_BG, AKAMAI_BLUE)
tb(slide, Inches(0.6), Inches(6.68), Inches(4.8), Inches(0.5),
   "NodeBalancer :80  \u2192  Demo App", sz=13, color=AKAMAI_BLUE, bold=True,
   align=PP_ALIGN.CENTER)

# User arrow
tb(slide, Inches(1.8), Inches(7.3), Inches(2.0), Inches(0.3),
   "\u2191  User Request", sz=10, color=CHARCOAL, align=PP_ALIGN.CENTER)

# Object Storage
card(slide, Inches(6.4), Inches(6.15), Inches(6.5), Inches(0.55), CARD_BG, AKAMAI_ORANGE)
tb(slide, Inches(6.6), Inches(6.18), Inches(6.1), Inches(0.5),
   "Object Storage (Chennai)  \u2014  CLIP + CLAP Model Weights",
   sz=13, color=AKAMAI_ORANGE, bold=True, align=PP_ALIGN.CENTER)

tb(slide, Inches(6.6), Inches(6.78), Inches(6.1), Inches(0.3),
   "S3-compatible  |  Cross-region pull from Mumbai compute", sz=10, color=GRAY,
   align=PP_ALIGN.CENTER)

# ==========================================
# Speaker Notes
# ==========================================
notes = slide.notes_slide
tf = notes.notes_text_frame

notes_text = """ARCHITECTURE OVERVIEW
This slide shows the complete system architecture for the CLIP/CLAP GPU inference demo running on Akamai Cloud Computing infrastructure.

LKE ENTERPRISE CLUSTER
- The entire cluster runs on LKE Enterprise (LKE-E), which provides a fully private VPC. The control plane, worker nodes, and all pod communications remain private with no public-facing endpoints on any node.
- A Cloud Firewall is automatically created and attached to all worker node VMs as part of every LKE-E cluster. VPC communication rules are pre-populated. Do NOT pass firewall_id when creating node pools via the API — this is a documented known issue that causes nodes to never reach ready state.
- The Kubernetes API is accessed via a Control Plane ACL — only traffic from explicitly allowed IP addresses can reach the control plane.

CPU NODE POOL (2x g6-standard-4)
- Ray Head runs with num-cpus: 0, meaning it acts purely as a scheduler and runs no inference workloads. This is a deliberate design choice to keep expensive GPU resources dedicated to model serving.
- KubeRay 1.6.0 operator manages the cluster via a consolidated RayService CRD — this single custom resource defines both the Ray cluster topology and the Serve deployment configuration. No standalone RayCluster resource is needed.
- The Demo App is a FastAPI backend serving a SvelteKit static frontend. It reverse-proxies Grafana at /grafana/ so dashboards are embedded in the UI without exposing Grafana as a separate service.
- Prometheus scrapes three targets: Ray Serve metrics from the head pod, FastAPI request metrics from the demo app, and NVIDIA DCGM GPU metrics from the GPU nodes.
- Grafana runs as a ClusterIP service (not LoadBalancer) and is accessed exclusively through the demo app's reverse proxy. This avoids a second external entry point.

GPU NODE POOL (3x g2-gpu-rtx4000a1-s — RTX 4000 Ada)
- Each node has 4 vCPU, 16GB RAM, and 1 GPU at $350/month.
- A NoSchedule taint (nvidia.com/gpu=true) prevents any non-inference workload from landing on GPU nodes. Only Ray workers carry the matching toleration.
- NVIDIA k8s-device-plugin v0.19.0 DaemonSet is required — LKE does not pre-install it. Without this, nodes have GPUs but the Kubernetes scheduler cannot see them (nvidia.com/gpu not advertised in node capacity).
- CLIP (openai/clip-vit-large-patch14) produces 768-dimensional embeddings for text and images. It runs with 2 replicas across 2 GPU nodes.
- CLAP (laion/clap-htsat-unfused) produces 512-dimensional embeddings for audio. It runs with 1 replica on the third GPU node.
- The serveConfigV2 in the RayService CRD defines two separate applications at /clip and /clap route prefixes. The head pod is excluded from serve traffic (excludeHeadPodFromServeSvc: true).

INIT CONTAINERS AND MODEL WEIGHT DISTRIBUTION
- Each GPU worker pod runs an init container before Ray starts. This container uses aws-cli to pull model weights from Linode Object Storage (S3-compatible) into an emptyDir volume mounted at /model-cache.
- Both CLIP and CLAP weights are pulled to every worker, since Ray Serve may schedule either model on any worker at deployment time.
- The emptyDir is ephemeral — weights are re-pulled on every pod restart. A PersistentVolume cache is a documented v2 optimization.
- Object Storage credentials are injected via a Kubernetes Secret created from Terraform outputs at deploy time.

CLOUD FIREWALL + NODEBALANCER
- The NodeBalancer is the ONLY externally accessible component. It is created automatically by Kubernetes when the demo app's LoadBalancer Service is applied.
- A separate Cloud Firewall (created by Terraform) is attached to the NodeBalancer by the deploy script after provisioning. This firewall restricts ports 80 and 443 to a configurable list of allowed IP CIDRs — no other traffic reaches the NodeBalancer.
- The NodeBalancer routes into the VPC using a /30 subnet from 10.254.0.0/22, configured via the service annotation service.beta.kubernetes.io/linode-loadbalancer-backend-ipv4-range.
- Traffic flow: User -> Cloud Firewall (IP check) -> NodeBalancer :80 -> FastAPI pod :8080 -> Ray Serve (cluster-internal) -> GPU worker.

OBJECT STORAGE (CHENNAI — in-maa-1)
- Model weights are stored in Linode Object Storage in Chennai because Mumbai (in-bom-2) does not offer Object Storage. This split-region design means compute runs in Mumbai and storage in Chennai.
- Cross-region model pulls add approximately 30-60 seconds to pod startup, but have zero impact on inference latency since models are cached locally on the pod after init.
- The S3 endpoint is derived from the Terraform bucket resource's s3_endpoint computed attribute — never constructed manually from the region string.
- Upload uses s5cmd for fast parallel transfer. The upload-models.sh script reads credentials from Terraform outputs.

DEPLOYMENT
- The entire stack deploys via four commands: terraform apply (infrastructure), build-images.sh (Docker images), upload-models.sh (model weights to OBJ storage), deploy.sh (K8s resources).
- The deploy script takes the Docker registry prefix as a required argument. Image references in YAML manifests use placeholders (WORKER_IMAGE_PLACEHOLDER, DEMO_IMAGE_PLACEHOLDER) that are injected via sed at deploy time.
- Region is fully parameterized — change region and obj_storage_region in terraform.tfvars to redeploy the entire stack in a different geography."""

for i, line in enumerate(notes_text.strip().split("\n")):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(10)
    p.font.name = FONT
    if line.isupper() and len(line) > 3:
        p.font.bold = True
        p.font.size = Pt(11)

# ==========================================
# SLIDE 2: Traditional flow architecture diagram
# ==========================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide2)

# Header
accent_bar(slide2, Inches(0), Inches(0), Inches(13.333), AKAMAI_BLUE)
tb(slide2, Inches(0.6), Inches(0.2), Inches(5), Inches(0.4),
   "AKAMAI CLOUD COMPUTING", sz=12, color=AKAMAI_BLUE, bold=True)
tb(slide2, Inches(0.6), Inches(0.55), Inches(10), Inches(0.6),
   "Request Flow + Infrastructure Topology", sz=26, color=NAVY, bold=True)
accent_bar(slide2, Inches(0.6), Inches(1.1), Inches(2.5))


def arrow_h(slide, x1, y, x2, color=GRAY):
    """Horizontal arrow using a thin rectangle + triangle."""
    # Line
    w = x2 - x1 - Inches(0.15)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x1, y, w, Pt(2))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    # Arrowhead
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, x2 - Inches(0.15), y - Inches(0.06), Inches(0.15), Inches(0.14))
    tri.rotation = 90.0
    tri.fill.solid()
    tri.fill.fore_color.rgb = color
    tri.line.fill.background()


def arrow_v(slide, x, y1, y2, color=GRAY):
    """Vertical downward arrow."""
    h = y2 - y1 - Inches(0.15)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y1, Pt(2), h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, x - Inches(0.06), y2 - Inches(0.15), Inches(0.14), Inches(0.15))
    tri.rotation = 180.0
    tri.fill.solid()
    tri.fill.fore_color.rgb = color
    tri.line.fill.background()


# Layout: left-to-right flow across the slide
# Row positions
ROW_MAIN = Inches(3.0)  # center of main flow row
NODE_W = Inches(2.0)
NODE_H = Inches(1.6)

# ---- User ----
card(slide2, Inches(0.3), Inches(2.6), Inches(1.6), Inches(1.8), CARD_BG, NAVY)
# Person icon (circle + body)
head = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), Inches(2.75), Inches(0.4), Inches(0.4))
head.fill.solid()
head.fill.fore_color.rgb = NAVY
head.line.fill.background()
tb(slide2, Inches(0.4), Inches(3.2), Inches(1.4), Inches(0.3),
   "User", sz=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
tb(slide2, Inches(0.4), Inches(3.5), Inches(1.4), Inches(0.6),
   "Text, Image,\nor Audio input", sz=10, color=CHARCOAL, align=PP_ALIGN.CENTER)

# Arrow: User → Firewall
arrow_h(slide2, Inches(1.95), ROW_MAIN + Inches(0.2), Inches(2.5), AKAMAI_BLUE)

# ---- Cloud Firewall ----
card(slide2, Inches(2.5), Inches(2.6), Inches(1.5), Inches(1.8), AKAMAI_BLUE, AKAMAI_BLUE)
tb(slide2, Inches(2.55), Inches(2.75), Inches(1.4), Inches(0.3),
   "Cloud", sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tb(slide2, Inches(2.55), Inches(3.0), Inches(1.4), Inches(0.3),
   "Firewall", sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tb(slide2, Inches(2.55), Inches(3.5), Inches(1.4), Inches(0.6),
   "IP allowlist\nPorts 80/443", sz=10, color=RGBColor(0xB0, 0xDD, 0xF0), align=PP_ALIGN.CENTER)

# Arrow: Firewall → NodeBalancer
arrow_h(slide2, Inches(4.05), ROW_MAIN + Inches(0.2), Inches(4.55), AKAMAI_BLUE)

# ---- NodeBalancer ----
card(slide2, Inches(4.55), Inches(2.6), Inches(1.5), Inches(1.8), CARD_BG, AKAMAI_BLUE)
tb(slide2, Inches(4.6), Inches(2.8), Inches(1.4), Inches(0.5),
   "Node\nBalancer", sz=13, color=AKAMAI_BLUE, bold=True, align=PP_ALIGN.CENTER)
tb(slide2, Inches(4.6), Inches(3.5), Inches(1.4), Inches(0.6),
   ":80 \u2192 :8080\nVPC backend", sz=10, color=CHARCOAL, align=PP_ALIGN.CENTER)

# Arrow: NB → FastAPI
arrow_h(slide2, Inches(6.1), ROW_MAIN + Inches(0.2), Inches(6.6), CHARCOAL)

# ---- VPC boundary starts ----
card(slide2, Inches(6.4), Inches(1.5), Inches(6.6), Inches(5.5), CARD_GRAY, BORDER_LIGHT)
tb(slide2, Inches(6.6), Inches(1.6), Inches(4), Inches(0.3),
   "LKE-E Private VPC", sz=11, color=CHARCOAL, bold=True, font=FONT_MONO)

# ---- FastAPI Demo App ----
card(slide2, Inches(6.6), Inches(2.6), Inches(1.7), Inches(1.8), CARD_BLUE, AKAMAI_BLUE)
tb(slide2, Inches(6.65), Inches(2.8), Inches(1.6), Inches(0.3),
   "FastAPI", sz=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
tb(slide2, Inches(6.65), Inches(3.1), Inches(1.6), Inches(0.3),
   "Demo App", sz=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
tb(slide2, Inches(6.65), Inches(3.5), Inches(1.6), Inches(0.6),
   "SvelteKit UI\n/embed/* routes\nGrafana proxy", sz=9, color=CHARCOAL, align=PP_ALIGN.CENTER)

# Arrow: FastAPI → Ray Serve
arrow_h(slide2, Inches(8.35), ROW_MAIN + Inches(0.2), Inches(8.85), CHARCOAL)

# ---- Ray Serve ----
card(slide2, Inches(8.85), Inches(2.2), Inches(1.7), Inches(2.6), CARD_BG, NAVY)
tb(slide2, Inches(8.9), Inches(2.3), Inches(1.6), Inches(0.3),
   "Ray Serve", sz=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
tb(slide2, Inches(8.9), Inches(2.6), Inches(1.6), Inches(0.25),
   "via RayService CRD", sz=8, color=GRAY, align=PP_ALIGN.CENTER, font=FONT_MONO)

# CLIP box
card(slide2, Inches(8.95), Inches(2.95), Inches(1.5), Inches(0.75), CARD_ORANGE, AKAMAI_ORANGE)
tb(slide2, Inches(9.0), Inches(2.98), Inches(1.4), Inches(0.25),
   "/clip", sz=11, color=AKAMAI_ORANGE, bold=True, align=PP_ALIGN.CENTER, font=FONT_MONO)
tb(slide2, Inches(9.0), Inches(3.25), Inches(1.4), Inches(0.35),
   "CLIP \u00d7 2 replicas\n768-dim embeddings", sz=8, color=CHARCOAL, align=PP_ALIGN.CENTER)

# CLAP box
card(slide2, Inches(8.95), Inches(3.8), Inches(1.5), Inches(0.75), CARD_ORANGE, AKAMAI_ORANGE)
tb(slide2, Inches(9.0), Inches(3.83), Inches(1.4), Inches(0.25),
   "/clap", sz=11, color=AKAMAI_ORANGE, bold=True, align=PP_ALIGN.CENTER, font=FONT_MONO)
tb(slide2, Inches(9.0), Inches(4.1), Inches(1.4), Inches(0.35),
   "CLAP \u00d7 1 replica\n512-dim embeddings", sz=8, color=CHARCOAL, align=PP_ALIGN.CENTER)

# Arrow: Ray Serve → GPU Workers
arrow_h(slide2, Inches(10.6), ROW_MAIN + Inches(0.2), Inches(11.1), AKAMAI_ORANGE)

# ---- GPU Workers ----
card(slide2, Inches(11.1), Inches(2.2), Inches(1.7), Inches(2.6), CARD_ORANGE, AKAMAI_ORANGE)
tb(slide2, Inches(11.15), Inches(2.3), Inches(1.6), Inches(0.3),
   "GPU Workers", sz=13, color=AKAMAI_ORANGE, bold=True, align=PP_ALIGN.CENTER)
tb(slide2, Inches(11.15), Inches(2.6), Inches(1.6), Inches(0.25),
   "3x RTX 4000 Ada", sz=9, color=GRAY, align=PP_ALIGN.CENTER, font=FONT_MONO)

for i, label in enumerate(["Worker 1\nCLIP", "Worker 2\nCLIP", "Worker 3\nCLAP"]):
    y = Inches(2.95 + i * 0.6)
    card(slide2, Inches(11.2), y, Inches(1.5), Inches(0.5), CARD_BG, BORDER_LIGHT)
    tb(slide2, Inches(11.25), y + Inches(0.02), Inches(1.4), Inches(0.45),
       label, sz=8, color=NAVY, align=PP_ALIGN.CENTER, font=FONT_MONO)

# ==========================================
# Object Storage — below GPU workers, with arrow up
# ==========================================
card(slide2, Inches(10.5), Inches(5.5), Inches(2.4), Inches(1.2), CARD_BG, AKAMAI_ORANGE)
tb(slide2, Inches(10.55), Inches(5.6), Inches(2.3), Inches(0.3),
   "Object Storage", sz=13, color=AKAMAI_ORANGE, bold=True, align=PP_ALIGN.CENTER)
tb(slide2, Inches(10.55), Inches(5.9), Inches(2.3), Inches(0.3),
   "Chennai (in-maa-1)", sz=10, color=CHARCOAL, align=PP_ALIGN.CENTER)
tb(slide2, Inches(10.55), Inches(6.2), Inches(2.3), Inches(0.3),
   "CLIP + CLAP weights", sz=9, color=GRAY, align=PP_ALIGN.CENTER, font=FONT_MONO)

# Arrow: OBJ Storage → GPU Workers (upward)
s = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.9), Inches(4.9), Pt(2), Inches(0.5))
s.fill.solid()
s.fill.fore_color.rgb = AKAMAI_ORANGE
s.line.fill.background()
tri = slide2.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.84), Inches(4.8), Inches(0.14), Inches(0.15))
tri.rotation = 0.0
tri.fill.solid()
tri.fill.fore_color.rgb = AKAMAI_ORANGE
tri.line.fill.background()
tb(slide2, Inches(10.6), Inches(5.15), Inches(1.2), Inches(0.25),
   "init pull", sz=8, color=AKAMAI_ORANGE, font=FONT_MONO)

# ==========================================
# Monitoring — below FastAPI
# ==========================================
card(slide2, Inches(6.6), Inches(5.5), Inches(2.8), Inches(1.2), CARD_BG, AKAMAI_BLUE)
tb(slide2, Inches(6.65), Inches(5.6), Inches(2.7), Inches(0.3),
   "Observability", sz=13, color=AKAMAI_BLUE, bold=True, align=PP_ALIGN.CENTER)
tb(slide2, Inches(6.65), Inches(5.9), Inches(2.7), Inches(0.6),
   "Prometheus + Grafana + DCGM\nGPU util, latency, throughput", sz=9, color=CHARCOAL,
   align=PP_ALIGN.CENTER)

# Arrow: Monitoring ↔ GPU (scrapes)
s = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5), Inches(6.1), Inches(0.9), Pt(2))
s.fill.solid()
s.fill.fore_color.rgb = GRAY
s.line.fill.background()
tb(slide2, Inches(9.55), Inches(5.85), Inches(0.8), Inches(0.2),
   "scrapes", sz=7, color=GRAY, align=PP_ALIGN.CENTER, font=FONT_MONO)

# Arrow: Monitoring → FastAPI (proxied via /grafana)
s = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(4.45), Pt(2), Inches(1.0))
s.fill.solid()
s.fill.fore_color.rgb = GRAY
s.line.fill.background()
tb(slide2, Inches(7.6), Inches(4.8), Inches(0.8), Inches(0.2),
   "/grafana", sz=7, color=GRAY, font=FONT_MONO)

# ==========================================
# Flow labels along the top
# ==========================================
flow_labels = [
    (Inches(0.6), "Internet"),
    (Inches(2.7), "Edge\nSecurity"),
    (Inches(4.7), "Load\nBalancing"),
    (Inches(6.8), "Application\nLayer"),
    (Inches(9.1), "Inference\nRouting"),
    (Inches(11.3), "GPU\nCompute"),
]
for x, label in flow_labels:
    tb(slide2, x, Inches(1.9), Inches(1.4), Inches(0.5),
       label, sz=9, color=GRAY, align=PP_ALIGN.CENTER)

# ==========================================
# Slide 2 Speaker Notes
# ==========================================
notes2 = slide2.notes_slide
tf2 = notes2.notes_text_frame

notes2_text = """REQUEST FLOW — LEFT TO RIGHT
Walk the audience through the request path from left to right. Each box is a component that the request passes through.

USER
- The user interacts with a web UI (SvelteKit) served by the demo app. They can submit text, upload an image, or upload audio.
- The UI calls the FastAPI backend's /embed/text, /embed/image, or /embed/audio endpoints.

CLOUD FIREWALL
- A Linode Cloud Firewall sits in front of the NodeBalancer. It is created by Terraform and attached to the NodeBalancer by the deploy script.
- Only traffic from explicitly allowed IP CIDRs (configured in terraform.tfvars as allowed_ips) can reach ports 80 and 443. All other inbound traffic is dropped before it reaches any hardware.
- This is separate from the auto-created LKE-E cluster firewall on the worker nodes — that one protects the nodes, this one protects the NodeBalancer.

NODEBALANCER
- The Linode NodeBalancer is the only externally addressable component. It is created automatically when Kubernetes applies the demo app's LoadBalancer Service.
- It terminates external HTTP on port 80 and forwards to the FastAPI pods on port 8080 inside the VPC.
- Backend routing uses a /30 subnet from the 10.254.0.0/22 range, configured via the service annotation service.beta.kubernetes.io/linode-loadbalancer-backend-ipv4-range.

LKE-E PRIVATE VPC
- Everything to the right of the NodeBalancer runs inside the LKE Enterprise private VPC. No nodes have public IPs. The control plane, worker nodes, and pod communications all remain private.
- The VPC boundary on the slide represents this isolation. Nothing inside is directly reachable from the internet.

FASTAPI DEMO APP
- The FastAPI application serves three roles: it serves the SvelteKit frontend as static files, it proxies embedding requests to Ray Serve, and it reverse-proxies Grafana dashboards at /grafana/ so they can be embedded in the UI without exposing Grafana externally.
- Two replicas run on CPU nodes (nodeSelector: workload=system).
- Prometheus scrapes the demo app's /metrics endpoint for HTTP request metrics (duration, status codes, in-flight requests).

RAY SERVE
- Ray Serve is managed by a single RayService CRD via KubeRay 1.6.0. The CRD defines two applications: /clip for CLIP embeddings and /clap for CLAP embeddings.
- The Ray head pod runs on a CPU node with num-cpus: 0 — it is a scheduler only, never runs inference. It is excluded from the serve traffic path (excludeHeadPodFromServeSvc: true).
- FastAPI calls the clip-clap-serve-svc ClusterIP service on port 8000. Ray routes the request to the appropriate deployment based on the route prefix.
- Prometheus scrapes the Ray head pod's :8080/metrics endpoint for serve-level metrics: request counts, latency histograms (p50/p99), replica health, error rates, and queue depth.

GPU WORKERS
- Three RTX 4000 Ada GPU nodes, each running one Ray worker with one GPU. CLIP occupies two workers, CLAP occupies one.
- A NoSchedule taint (nvidia.com/gpu=true) on the GPU node pool ensures only Ray workers land on these nodes.
- The NVIDIA k8s-device-plugin DaemonSet (v0.19.0) must be installed — LKE does not pre-install it. Without it, the scheduler cannot see the GPUs.
- DCGM exporter runs as a DaemonSet on GPU nodes and exposes GPU utilization, memory, temperature, and power draw metrics to Prometheus.

OBJECT STORAGE (CHENNAI)
- Model weights (CLIP: openai/clip-vit-large-patch14, CLAP: laion/clap-htsat-unfused) are stored in Linode Object Storage in Chennai (in-maa-1) because Mumbai (in-bom-2) does not offer Object Storage.
- Each GPU worker pod runs an init container that pulls both model sets from Object Storage using aws-cli with S3-compatible credentials before Ray starts.
- After init, models are cached in an emptyDir volume at /model-cache. All inference is GPU-local — zero network calls during request processing.
- The cross-region pull (Mumbai to Chennai) adds ~30-60 seconds to pod startup but has no impact on inference latency.

OBSERVABILITY
- Prometheus scrapes three targets: Ray head pod (serve metrics), demo app pods (HTTP metrics), and DCGM exporter (GPU hardware metrics).
- Grafana runs as a ClusterIP service inside the VPC. It is NOT exposed via LoadBalancer. Instead, the demo app reverse-proxies it at /grafana/, keeping it accessible in the UI without a second external entry point.
- Two provisioned dashboards: Ray Serve (requests/sec, p50/p99 latency, replica count, errors) and GPU Utilization (util%, memory, temperature, power per card)."""

for i, line in enumerate(notes2_text.strip().split("\n")):
    if i == 0:
        p = tf2.paragraphs[0]
    else:
        p = tf2.add_paragraph()
    p.text = line
    p.font.size = Pt(10)
    p.font.name = FONT
    if line.isupper() and len(line) > 3:
        p.font.bold = True
        p.font.size = Pt(11)

# ==========================================
out = "/home/morduun/projects/akamai/ray-dist-gpu/CLIP-CLAP-Architecture.pptx"
prs.save(out)
print(f"Saved to {out}")
