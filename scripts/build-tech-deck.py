#!/usr/bin/env python3
"""Generate the technical deep-dive deck for the CLIP/CLAP inference network."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# -- Brand palette --
BG_DARK = RGBColor(0x0A, 0x0A, 0x0F)
BG_CARD = RGBColor(0x12, 0x12, 0x1A)
BG_ELEVATED = RGBColor(0x1A, 0x1A, 0x26)
BG_CODE = RGBColor(0x0E, 0x0E, 0x16)
ACCENT = RGBColor(0x00, 0xE5, 0xC8)
ACCENT_AMBER = RGBColor(0xF0, 0xA0, 0x30)
ACCENT_BLUE = RGBColor(0x60, 0x90, 0xFF)
TEXT_PRIMARY = RGBColor(0xE8, 0xE8, 0xF0)
TEXT_SECONDARY = RGBColor(0x88, 0x88, 0xA0)
TEXT_MUTED = RGBColor(0x55, 0x55, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0x2A, 0x2A, 0x3E)
RED_SOFT = RGBColor(0xF0, 0x50, 0x50)
GREEN = RGBColor(0x40, 0xE0, 0x70)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_MONO = "Consolas"
FONT_DISPLAY = "Calibri"


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_DISPLAY):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=TEXT_PRIMARY, bullet_color=ACCENT):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run_b = p.add_run()
        run_b.text = "\u25B8  "
        run_b.font.size = Pt(font_size)
        run_b.font.color.rgb = bullet_color
        run_b.font.name = FONT_DISPLAY
        run_t = p.add_run()
        run_t.text = item
        run_t.font.size = Pt(font_size)
        run_t.font.color.rgb = color
        run_t.font.name = FONT_DISPLAY
    return tf


def add_code_block(slide, left, top, width, height, code, font_size=10):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CODE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.5)
    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.strip().split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = TEXT_PRIMARY
        run.font.name = FONT_MONO
    return tf


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def section_header(slide, label, title, subtitle=None):
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
                label, font_size=12, color=ACCENT, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.5))
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.8),
                title, font_size=34, color=TEXT_PRIMARY, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(1.85), Inches(11), Inches(0.5),
                    subtitle, font_size=17, color=TEXT_SECONDARY)


def kv_row(slide, left, y, key, val, key_w=Inches(2.5), val_w=Inches(4),
           key_color=ACCENT, val_color=TEXT_PRIMARY, font_size=13):
    add_textbox(slide, left, y, key_w, Inches(0.3), key,
                font_size=font_size, color=key_color, bold=True, font_name=FONT_MONO)
    add_textbox(slide, left + key_w, y, val_w, Inches(0.3), val,
                font_size=font_size, color=val_color)


# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ----------------------------------------------------------
# SLIDE 1: Title
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
            "TECHNICAL DEEP DIVE", font_size=14, color=ACCENT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.8), Inches(2))
add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(1.2),
            "GPU Inference on Akamai LKE\nwith KubeRay + Ray Serve", font_size=44, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(3.8), Inches(8), Inches(0.5),
            "CLIP/CLAP Multimodal Embedding Service — Architecture, Implementation, and Lessons Learned",
            font_size=18, color=TEXT_SECONDARY)

# Version badges
versions = [
    ("Ray 2.54.0", ACCENT), ("KubeRay 1.6.0", ACCENT),
    ("K8s 1.35", ACCENT_BLUE), ("Terraform ~> 3.8", ACCENT_BLUE),
    ("Python 3.11", ACCENT_AMBER), ("SvelteKit 2", ACCENT_AMBER),
]
for i, (v, c) in enumerate(versions):
    x = Inches(0.8 + i * 1.8)
    add_card(slide, x, Inches(5.2), Inches(1.6), Inches(0.4), BG_ELEVATED)
    add_textbox(slide, x + Inches(0.1), Inches(5.22), Inches(1.4), Inches(0.35),
                v, font_size=11, color=c, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_MONO)

# ----------------------------------------------------------
# SLIDE 2: Stack Overview
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
section_header(slide, "ARCHITECTURE", "Full Stack Overview")

layers = [
    ("LAYER 5: PRESENTATION", "SvelteKit static SPA \u2192 FastAPI static mount \u2192 NodeBalancer :80",
     [("Frontend", "SvelteKit 2 + adapter-static, dark industrial aesthetic"),
      ("Backend", "FastAPI 0.111+, proxies to Ray Serve, serves static build"),
      ("Ingress", "Linode NodeBalancer via LoadBalancer Service")], ACCENT),

    ("LAYER 4: INFERENCE", "Ray Serve deployments on GPU workers, managed by KubeRay RayService CRD",
     [("CLIP", "openai/clip-vit-large-patch14 \u2014 768-dim text+image embeddings, 2 replicas"),
      ("CLAP", "laion/clap-htsat-unfused \u2014 512-dim audio embeddings, 1 replica"),
      ("Routing", "/clip and /clap route prefixes, head excluded from serve traffic")], ACCENT_AMBER),

    ("LAYER 3: ORCHESTRATION", "KubeRay operator + NVIDIA device plugin on LKE",
     [("KubeRay 1.6.0", "RayService CRD manages cluster lifecycle + serve deployments"),
      ("Device Plugin", "v0.19.0 DaemonSet exposes nvidia.com/gpu to scheduler"),
      ("GPU Taints", "NoSchedule taint on GPU pool, only Ray workers tolerate")], ACCENT_BLUE),

    ("LAYER 2: INFRASTRUCTURE", "Terraform-managed LKE cluster with VPC + firewall",
     [("Compute", "in-bom-2 (Mumbai): 2x CPU g6-standard-4 + 3x GPU g2-gpu-rtx4000a1-s"),
      ("Storage", "in-maa-1 (Chennai): S3-compatible Object Storage for model weights"),
      ("Network", "VPC + private subnet + Cloud Firewall, operator-restricted access")], TEXT_SECONDARY),
]

for i, (layer_name, desc, items, color) in enumerate(layers):
    y = Inches(2.2 + i * 1.25)
    add_card(slide, Inches(0.6), y, Inches(12.1), Inches(1.1))
    add_textbox(slide, Inches(0.9), y + Inches(0.05), Inches(3), Inches(0.3),
                layer_name, font_size=10, color=color, bold=True, font_name=FONT_MONO)
    add_textbox(slide, Inches(0.9), y + Inches(0.3), Inches(3.5), Inches(0.3),
                desc, font_size=10, color=TEXT_MUTED)
    for j, (k, v) in enumerate(items):
        x = Inches(4.8 + j * 2.8)
        add_textbox(slide, x, y + Inches(0.05), Inches(2.6), Inches(0.3),
                    k, font_size=11, color=color, bold=True, font_name=FONT_MONO)
        add_textbox(slide, x, y + Inches(0.35), Inches(2.6), Inches(0.7),
                    v, font_size=10, color=TEXT_SECONDARY)

# ----------------------------------------------------------
# SLIDE 3: Terraform Module Design
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
section_header(slide, "INFRASTRUCTURE AS CODE", "Terraform Module Architecture",
               "3 modules, 15 files, provider ~> 3.8 with v4beta API for VPC support")

# Module diagram
modules_data = [
    ("networking/", "VPC \u2192 Subnet \u2192 Firewall",
     ["linode_vpc \u2014 region-scoped private network",
      "linode_vpc_subnet \u2014 10.0.0.0/16 CIDR",
      "linode_firewall \u2014 per-port ACLs, operator_ips gated",
      "Outputs: vpc_id, subnet_id, firewall_id"],
     ACCENT),
    ("lke/", "Cluster \u2192 Node Pools \u2192 Kubeconfig",
     ["linode_lke_cluster \u2014 VPC+subnet at cluster level",
      "CPU pool: 2x g6-standard-4, labels: workload=system",
      "GPU pool: 3x g2-gpu-rtx4000a1-s, taint: nvidia.com/gpu",
      "firewall_id per pool (provider >= 3.8 required)",
      "local_sensitive_file \u2014 kubeconfig to disk"],
     ACCENT_AMBER),
    ("storage/", "Bucket \u2192 Access Keys",
     ["linode_object_storage_bucket \u2014 private ACL",
      "Reader key: init containers in Ray workers",
      "Writer key: upload-models.sh one-time use",
      "Endpoint from bucket.s3_endpoint (not constructed)"],
     ACCENT_BLUE),
]

for i, (name, desc, details, color) in enumerate(modules_data):
    y_base = Inches(2.4 + i * 1.7)
    add_card(slide, Inches(0.8), y_base, Inches(11.7), Inches(1.5))
    add_textbox(slide, Inches(1.0), y_base + Inches(0.05), Inches(2), Inches(0.3),
                f"modules/{name}", font_size=14, color=color, bold=True, font_name=FONT_MONO)
    add_textbox(slide, Inches(1.0), y_base + Inches(0.35), Inches(3), Inches(0.3),
                desc, font_size=12, color=TEXT_SECONDARY)
    for j, d in enumerate(details):
        add_textbox(slide, Inches(4.5), y_base + Inches(0.05 + j * 0.3), Inches(7.5), Inches(0.3),
                    d, font_size=11, color=TEXT_PRIMARY, font_name=FONT_MONO)

# ----------------------------------------------------------
# SLIDE 4: VPC + Firewall Detail
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
section_header(slide, "NETWORKING", "VPC + Cloud Firewall Design",
               "LKE requires v4beta API for VPC attachment. Firewall binds per-pool, not per-cluster.")

# Key insight
add_card(slide, Inches(0.8), Inches(2.4), Inches(5.5), Inches(1.2), BG_ELEVATED)
add_textbox(slide, Inches(1.0), Inches(2.5), Inches(5), Inches(0.3),
            "CRITICAL DISCOVERY", font_size=11, color=RED_SOFT, bold=True)
add_textbox(slide, Inches(1.0), Inches(2.85), Inches(5), Inches(0.7),
            "provider \"linode\" requires api_version = \"v4beta\" for\nvpc_id and subnet_id on linode_lke_cluster.\nProvider must be >= 3.8 for firewall_id on pools.",
            font_size=13, color=TEXT_PRIMARY, font_name=FONT_MONO)

# Firewall rules table
add_textbox(slide, Inches(0.8), Inches(3.9), Inches(5), Inches(0.3),
            "Firewall Rules", font_size=18, color=TEXT_PRIMARY, bold=True)

rules = [
    ("80, 443", "operator_ips", "Demo app via NodeBalancer"),
    ("6443", "operator_ips", "Kubernetes API"),
    ("3000", "operator_ips", "Grafana dashboard"),
    ("1-65535", "VPC CIDR", "Intra-cluster (TCP+UDP)"),
    ("*", "*", "DROP (default inbound policy)"),
]
for i, (ports, source, desc) in enumerate(rules):
    y = Inches(4.3 + i * 0.4)
    color = RED_SOFT if "DROP" in desc else TEXT_PRIMARY
    add_textbox(slide, Inches(0.8), y, Inches(1.2), Inches(0.3),
                ports, font_size=11, color=ACCENT, font_name=FONT_MONO)
    add_textbox(slide, Inches(2.1), y, Inches(1.5), Inches(0.3),
                source, font_size=11, color=TEXT_SECONDARY, font_name=FONT_MONO)
    add_textbox(slide, Inches(3.8), y, Inches(2.5), Inches(0.3),
                desc, font_size=11, color=color)

# Right: provider config
add_code_block(slide, Inches(7.0), Inches(2.4), Inches(5.5), Inches(2.2),
"""provider "linode" {
  token       = var.linode_token
  api_version = "v4beta"  # Required for VPC on LKE
}

# Each module needs required_providers:
terraform {
  required_providers {
    linode = {
      source  = "linode/linode"
      version = ">= 3.8"  # firewall_id on pools
    }
  }
}""", font_size=11)

add_code_block(slide, Inches(7.0), Inches(4.9), Inches(5.5), Inches(1.8),
"""# Firewall attaches PER POOL, not per cluster
pool {
  type        = var.gpu_node_type
  count       = var.gpu_node_count
  firewall_id = var.firewall_id  # <-- here

  taint {
    key    = "nvidia.com/gpu"
    value  = "true"
    effect = "NoSchedule"
  }
}""", font_size=11)

# ----------------------------------------------------------
# SLIDE 5: Split-Region Architecture
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
section_header(slide, "REGION TOPOLOGY", "Split-Region: Compute + Storage",
               "No single India region offers both RTX 4000 Ada GPUs and Object Storage.")

# Two region cards
for i, (region, city, components, color) in enumerate([
    ("in-bom-2", "Mumbai", ["LKE cluster (CPU + GPU pools)", "VPC + Cloud Firewall",
     "KubeRay + Ray Serve", "NVIDIA device plugin", "Prometheus + Grafana + DCGM", "Demo app"], ACCENT),
    ("in-maa-1", "Chennai", ["Object Storage bucket", "CLIP model weights (~1.7GB)",
     "CLAP model weights (~1.5GB)", "S3-compatible endpoint",
     "Accessed by init containers at pod startup", "Cross-region pull: ~30-60s per pod"], ACCENT_AMBER),
]):
    x = Inches(0.8 + i * 6.3)
    add_card(slide, x, Inches(2.4), Inches(5.8), Inches(4.2))
    add_textbox(slide, x + Inches(0.3), Inches(2.5), Inches(3), Inches(0.4),
                region, font_size=22, color=color, bold=True, font_name=FONT_MONO)
    add_textbox(slide, x + Inches(0.3), Inches(2.9), Inches(3), Inches(0.3),
                city, font_size=16, color=TEXT_SECONDARY)
    add_bullet_list(slide, x + Inches(0.3), Inches(3.4), Inches(5), Inches(3.0),
                    components, font_size=14, bullet_color=color)

# Variable config
add_code_block(slide, Inches(3.0), Inches(6.2), Inches(7.3), Inches(0.7),
"""# terraform.tfvars — two variables, full region control
region             = \"in-bom-2\"    # compute
obj_storage_region = \"in-maa-1\"    # storage""", font_size=12)

# ----------------------------------------------------------
# SLIDE 6: RayService CRD Deep Dive
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
section_header(slide, "MODEL SERVING", "RayService CRD — Consolidated Design",
               "Single CRD manages both the Ray cluster and Serve deployments. No standalone RayCluster needed.")

# Left: Why consolidated
add_textbox(slide, Inches(0.8), Inches(2.4), Inches(5.5), Inches(0.3),
            "Why RayService over RayCluster + manual serve?", font_size=16, color=TEXT_PRIMARY, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(2.8), Inches(5.5), Inches(2.5), [
    "Declarative serve config — no 'serve run' in a script",
    "Rolling upgrades via upgradeStrategy: NewCluster",
    "Health-aware traffic switching — old cluster stays live until new is healthy",
    "Single kubectl apply deploys everything",
    "excludeHeadPodFromServeSvc: true keeps head off GPU traffic path",
], font_size=14)

# Key config decisions
add_textbox(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(0.3),
            "Key Configuration Decisions", font_size=16, color=TEXT_PRIMARY, bold=True)

decisions = [
    ("serveConfigV2", "YAML string (not structured), posted to Ray REST API"),
    ("proxy_location", "EveryNode \u2014 proxy on each worker for local routing"),
    ("num-cpus: \"0\"", "Head does NO compute, scheduler only"),
    ("Two applications", "/clip (2 replicas) + /clap (1 replica) = 3 GPUs total"),
    ("preStop hook", "ray stop for graceful drain during upgrades"),
]
for i, (k, v) in enumerate(decisions):
    kv_row(slide, Inches(0.8), Inches(5.35 + i * 0.35), k, v, font_size=12)

# Right: serveConfigV2 snippet
add_code_block(slide, Inches(7.0), Inches(2.4), Inches(5.5), Inches(4.5),
"""serveConfigV2: |
  proxy_location: EveryNode
  http_options:
    host: 0.0.0.0
    port: 8000
    request_timeout_s: 120
  applications:
    - name: clip
      import_path: clip_deployment:app
      route_prefix: /clip
      deployments:
        - name: CLIPEmbedder
          num_replicas: 2
          ray_actor_options:
            num_gpus: 1
            num_cpus: 1
    - name: clap
      import_path: clap_deployment:app
      route_prefix: /clap
      deployments:
        - name: CLAPEmbedder
          num_replicas: 1
          ray_actor_options:
            num_gpus: 1""", font_size=10)

# ----------------------------------------------------------
# SLIDE 7: GPU Scheduling
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
section_header(slide, "GPU SCHEDULING", "NVIDIA Device Plugin + Taint/Toleration Strategy",
               "LKE GPU nodes ship without a device plugin. Workers can't schedule until GPUs are advertised.")

# The problem
add_card(slide, Inches(0.8), Inches(2.4), Inches(5.5), Inches(2.0), BG_ELEVATED)
add_textbox(slide, Inches(1.0), Inches(2.5), Inches(4), Inches(0.3),
            "THE PROBLEM", font_size=12, color=RED_SOFT, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(2.9), Inches(5), Inches(1.3), [
    "GPU nodes exist and have NVIDIA hardware",
    "But no DaemonSet exposes nvidia.com/gpu to K8s scheduler",
    "Workers request nvidia.com/gpu: 1 \u2192 Insufficient resources",
    "Pods stuck in Pending indefinitely",
], font_size=13, bullet_color=RED_SOFT)

# The solution
add_card(slide, Inches(0.8), Inches(4.7), Inches(5.5), Inches(2.2), BG_ELEVATED)
add_textbox(slide, Inches(1.0), Inches(4.8), Inches(4), Inches(0.3),
            "THE FIX", font_size=12, color=GREEN, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(5.15), Inches(5), Inches(1.5), [
    "Install NVIDIA k8s-device-plugin v0.19.0 DaemonSet",
    "Runs on ALL nodes \u2014 CrashLoops on CPU nodes (expected/harmless)",
    "GPU nodes report nvidia.com/gpu: 1 in capacity",
    "Workers schedule immediately, init containers pull weights",
    "Added to deploy.sh Phase 5 (before KubeRay)",
], font_size=13, bullet_color=GREEN)

# Right: scheduling flow
add_textbox(slide, Inches(7.0), Inches(2.4), Inches(5), Inches(0.3),
            "Scheduling Chain", font_size=18, color=TEXT_PRIMARY, bold=True)

chain = [
    ("1. Node Labels", "workload: inference\nnvidia.com/gpu: true", "Set by LKE pool config"),
    ("2. GPU Taint", "nvidia.com/gpu=true:NoSchedule", "Repels non-GPU workloads"),
    ("3. Device Plugin", "Exposes nvidia.com/gpu: 1\nas allocatable resource", "DaemonSet on GPU nodes"),
    ("4. Worker Toleration", "Tolerates nvidia.com/gpu taint", "In RayService workerGroupSpecs"),
    ("5. Resource Request", "nvidia.com/gpu: 1 in\nrequests AND limits", "K8s GPU scheduling requires both"),
]
for i, (step, detail, note) in enumerate(chain):
    y = Inches(2.9 + i * 0.9)
    add_card(slide, Inches(7.0), y, Inches(5.5), Inches(0.75))
    add_textbox(slide, Inches(7.2), y + Inches(0.02), Inches(2.0), Inches(0.3),
                step, font_size=12, color=ACCENT, bold=True, font_name=FONT_MONO)
    add_textbox(slide, Inches(7.2), y + Inches(0.3), Inches(2.5), Inches(0.4),
                detail, font_size=10, color=TEXT_PRIMARY, font_name=FONT_MONO)
    add_textbox(slide, Inches(9.8), y + Inches(0.15), Inches(2.5), Inches(0.4),
                note, font_size=10, color=TEXT_MUTED)

# ----------------------------------------------------------
# SLIDE 8: Model Weight Distribution
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
section_header(slide, "DATA FLOW", "Model Weight Distribution Pipeline",
               "Weights flow: HuggingFace \u2192 Object Storage \u2192 init container \u2192 emptyDir \u2192 Ray worker")

# Pipeline stages
stages = [
    ("UPLOAD (one-time)", "scripts/upload-models.sh",
     ["hf download \u2192 local tmpdir", "s5cmd sync \u2192 Object Storage",
      "Uses writer key from Terraform outputs", "S3_ENDPOINT_URL set from tf output obj_endpoint"],
     ACCENT),
    ("INIT (per pod startup)", "initContainers: model-puller",
     ["aws s3 sync from OBJ Storage \u2192 /model-cache", "Runs before Ray container starts",
      "Uses reader key from K8s Secret", "Cross-region Mumbai\u2192Chennai: ~30-60s for ~4GB"],
     ACCENT_AMBER),
    ("SERVE (runtime)", "Ray worker container",
     ["MODEL_CACHE_PATH=/model-cache env var", "transformers from_pretrained(path)",
      "Model loaded to GPU VRAM once at actor init", "Zero network calls during inference"],
     GREEN),
]

for i, (title, source, details, color) in enumerate(stages):
    x = Inches(0.8 + i * 4.1)
    add_card(slide, x, Inches(2.4), Inches(3.8), Inches(4.0))
    add_textbox(slide, x + Inches(0.2), Inches(2.5), Inches(3.4), Inches(0.3),
                title, font_size=12, color=color, bold=True)
    add_textbox(slide, x + Inches(0.2), Inches(2.8), Inches(3.4), Inches(0.3),
                source, font_size=11, color=TEXT_SECONDARY, font_name=FONT_MONO)
    add_bullet_list(slide, x + Inches(0.2), Inches(3.2), Inches(3.4), Inches(2.5),
                    details, font_size=12, color=TEXT_PRIMARY, bullet_color=color)

# Volume config
add_code_block(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8),
"""volumes:
  - name: model-cache
    emptyDir: { sizeLimit: 10Gi }    # Ephemeral — re-pulled on pod restart. PVC cache is a v2 optimization.""",
               font_size=11)

# ----------------------------------------------------------
# SLIDE 9: Deploy Pipeline
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
section_header(slide, "DEPLOYMENT", "9-Phase Deploy Pipeline",
               "scripts/deploy.sh <registry> [tag] — no hardcoded values, everything injected at runtime")

phases = [
    ("1", "Read TF Outputs", "OBJ keys, endpoint,\nbucket, Grafana pwd"),
    ("2", "Kubeconfig", "Extract from TF state,\nwrite to ./kubeconfig"),
    ("3", "Namespaces", "ray-system +\nmonitoring"),
    ("4", "Secrets/Config", "OBJ reader key +\ndynamic ConfigMap"),
    ("5", "Device Plugin", "NVIDIA k8s-device-plugin\nv0.19.0 DaemonSet"),
    ("6", "KubeRay", "Helm 1.6.0\nray-system namespace"),
    ("7", "RayService", "sed injects image,\nkubectl apply -f -"),
    ("8", "Monitoring", "Prometheus + DCGM +\nGrafana dashboards"),
    ("9", "Demo App", "sed injects image,\nrollout status wait"),
]

for i, (num, title, desc) in enumerate(phases):
    col = i % 5
    row = i // 5
    x = Inches(0.5 + col * 2.5)
    y = Inches(2.4 + row * 2.4)
    add_card(slide, x, y, Inches(2.3), Inches(2.0))
    # Number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.15),
                                     Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    add_textbox(slide, x + Inches(0.15), y + Inches(0.17), Inches(0.4), Inches(0.35),
                num, font_size=14, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.65), y + Inches(0.18), Inches(1.5), Inches(0.35),
                title, font_size=13, color=TEXT_PRIMARY, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.7), Inches(2.0), Inches(1.0),
                desc, font_size=11, color=TEXT_SECONDARY, font_name=FONT_MONO)

# Key technique callout
add_card(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), BG_ELEVATED)
add_textbox(slide, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.6),
            "Key technique: YAML manifests contain WORKER_IMAGE_PLACEHOLDER / DEMO_IMAGE_PLACEHOLDER. "
            "deploy.sh uses sed to inject the registry + tag at apply time. No hardcoded image references anywhere.",
            font_size=12, color=ACCENT, font_name=FONT_MONO)

# ----------------------------------------------------------
# SLIDE 10: Monitoring Architecture
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
section_header(slide, "OBSERVABILITY", "Monitoring Stack Design",
               "kube-prometheus-stack + DCGM exporter + custom Grafana dashboards")

# Three scrape targets
targets = [
    ("ray-serve", "Ray head pod :8080/metrics",
     ["ray_serve_num_requests_total", "ray_serve_request_latency_ms_bucket",
      "ray_serve_deployment_replica_healthy", "ray_serve_num_errors_total",
      "ray_serve_num_ongoing_requests"],
     "kubernetes_sd_configs with\npod label: ray.io/node-type=head", ACCENT),

    ("demo-app", "FastAPI pods :8080/metrics",
     ["http_request_duration_seconds", "http_requests_total",
      "http_requests_in_progress", "process_resident_memory_bytes"],
     "kubernetes_sd_configs with\npod label: app=demo-app", ACCENT_BLUE),

    ("dcgm-exporter", "GPU node DaemonSet :9400",
     ["DCGM_FI_DEV_GPU_UTIL", "DCGM_FI_DEV_FB_USED / FB_TOTAL",
      "DCGM_FI_DEV_GPU_TEMP", "DCGM_FI_DEV_POWER_USAGE"],
     "ServiceMonitor with label:\nrelease=kube-prometheus-stack", ACCENT_AMBER),
]

for i, (job, endpoint, metrics, discovery, color) in enumerate(targets):
    x = Inches(0.8 + i * 4.1)
    add_card(slide, x, Inches(2.4), Inches(3.8), Inches(4.5))
    add_textbox(slide, x + Inches(0.2), Inches(2.5), Inches(3.4), Inches(0.3),
                job, font_size=14, color=color, bold=True, font_name=FONT_MONO)
    add_textbox(slide, x + Inches(0.2), Inches(2.85), Inches(3.4), Inches(0.3),
                endpoint, font_size=11, color=TEXT_SECONDARY, font_name=FONT_MONO)
    add_accent_line(slide, x + Inches(0.2), Inches(3.2), Inches(1.5))

    add_textbox(slide, x + Inches(0.2), Inches(3.35), Inches(3), Inches(0.25),
                "Metrics:", font_size=10, color=TEXT_MUTED)
    for j, m in enumerate(metrics):
        add_textbox(slide, x + Inches(0.2), Inches(3.6 + j * 0.3), Inches(3.4), Inches(0.3),
                    m, font_size=10, color=TEXT_PRIMARY, font_name=FONT_MONO)

    add_textbox(slide, x + Inches(0.2), Inches(5.3), Inches(3.4), Inches(0.5),
                discovery, font_size=10, color=TEXT_MUTED, font_name=FONT_MONO)

# DCGM OOM note
add_card(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.7), BG_ELEVATED)
add_textbox(slide, Inches(1.0), Inches(6.55), Inches(1.5), Inches(0.3),
            "GOTCHA", font_size=11, color=RED_SOFT, bold=True)
add_textbox(slide, Inches(2.5), Inches(6.55), Inches(9.5), Inches(0.6),
            "DCGM exporter 4.8.x requires 512Mi+ memory (was OOMKilled at 256Mi). "
            "Production values: requests 512Mi, limits 1Gi. Also: the chart moved from "
            "nvidia/dcgm-exporter to https://nvidia.github.io/dcgm-exporter/helm-charts",
            font_size=12, color=TEXT_SECONDARY)

# ----------------------------------------------------------
# SLIDE 11: FastAPI + Frontend Architecture
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
section_header(slide, "APPLICATION LAYER", "FastAPI Backend + SvelteKit Frontend",
               "Single container serves API routes and static frontend. Multi-stage Docker build.")

# API routes
add_textbox(slide, Inches(0.8), Inches(2.4), Inches(5), Inches(0.3),
            "API Surface", font_size=18, color=TEXT_PRIMARY, bold=True)

routes = [
    ("POST /embed/text", "Text \u2192 CLIP \u2192 768-dim embedding", "Proxy to Ray /clip"),
    ("POST /embed/image", "Multipart image \u2192 CLIP", "Base64 encode, proxy to Ray /clip"),
    ("POST /embed/audio", "Multipart audio \u2192 CLAP", "Base64 encode, proxy to Ray /clap"),
    ("POST /similarity", "Two vectors \u2192 cosine score", "CPU-only numpy, no GPU"),
    ("GET /health", "Liveness + Ray reachability", "Probes Ray /-/healthz"),
    ("GET /config", "Runtime config for frontend", "Grafana URL, Ray Serve URL"),
    ("GET /metrics", "Prometheus exposition", "prometheus-fastapi-instrumentator"),
    ("GET /*", "SvelteKit static files", "app.mount(\"/\", StaticFiles(...))"),
]

for i, (route, desc, impl) in enumerate(routes):
    y = Inches(2.8 + i * 0.38)
    bg = BG_CARD if i % 2 == 0 else BG_DARK
    add_card(slide, Inches(0.8), y, Inches(5.5), Inches(0.35), bg)
    add_textbox(slide, Inches(0.9), y + Inches(0.02), Inches(1.8), Inches(0.3),
                route, font_size=10, color=ACCENT, font_name=FONT_MONO)
    add_textbox(slide, Inches(2.8), y + Inches(0.02), Inches(1.8), Inches(0.3),
                desc, font_size=10, color=TEXT_PRIMARY)
    add_textbox(slide, Inches(4.7), y + Inches(0.02), Inches(1.5), Inches(0.3),
                impl, font_size=9, color=TEXT_MUTED, font_name=FONT_MONO)

# Right: Frontend
add_textbox(slide, Inches(7.0), Inches(2.4), Inches(5), Inches(0.3),
            "Frontend Components", font_size=18, color=TEXT_PRIMARY, bold=True)

components = [
    ("EmbedInput.svelte", "Tabbed input (text/image/audio), dispatches events"),
    ("EmbedVisualizer.svelte", "Waveform + PCA scatter plot, animated on result"),
    ("SimilarityDemo.svelte", "Two-text compare, tweened score animation"),
    ("MetricsPanel.svelte", "Grafana iframe, connection status indicator"),
    ("ArchDiagram.svelte", "Inline SVG, no external dependencies"),
    ("api.ts", "Typed fetch wrappers for all backend routes"),
]

for i, (comp, desc) in enumerate(components):
    y = Inches(2.8 + i * 0.6)
    add_card(slide, Inches(7.0), y, Inches(5.5), Inches(0.5))
    add_textbox(slide, Inches(7.2), y + Inches(0.02), Inches(2.5), Inches(0.25),
                comp, font_size=11, color=ACCENT_AMBER, bold=True, font_name=FONT_MONO)
    add_textbox(slide, Inches(7.2), y + Inches(0.25), Inches(5), Inches(0.25),
                desc, font_size=10, color=TEXT_SECONDARY)

# Docker build
add_code_block(slide, Inches(7.0), Inches(6.2), Inches(5.5), Inches(1.0),
"""# Multi-stage: node builds frontend, python serves
FROM node:20-slim AS frontend-builder  # npm ci + build
FROM python:3.11-slim                  # pip install + copy build/
# Build context: app/ (not app/backend/)
# docker build -f backend/Dockerfile -t demo app/""", font_size=10)

# ----------------------------------------------------------
# SLIDE 12: Lessons Learned
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
section_header(slide, "LESSONS LEARNED", "Production Gotchas on Akamai LKE + GPU")

lessons = [
    ("Region Capability Gaps",
     "Not all regions have all features. us-east lacks VPCs. in-maa lacks RTX 4000 Ada GPUs. "
     "in-bom-2 lacks Object Storage. We ended up split-region: compute in Mumbai, storage in Chennai.",
     "Always validate region capabilities via API before committing. "
     "Design for split-region from day one.",
     RED_SOFT),

    ("NVIDIA Device Plugin Not Pre-Installed",
     "LKE GPU nodes have NVIDIA drivers but no k8s-device-plugin. "
     "Pods requesting nvidia.com/gpu get stuck in Pending with 'Insufficient' errors.",
     "Install k8s-device-plugin DaemonSet before any GPU workloads. "
     "Expect CrashLoopBackOff on CPU nodes (harmless).",
     ACCENT_AMBER),

    ("DCGM Exporter Memory Requirements",
     "DCGM 4.8.x OOMKills at 256Mi. Process initializes fine, then dies silently "
     "during first metric collection. Exit code 137 confirms OOM.",
     "Set limits to 1Gi. Also: chart repo moved from helm.ngc.nvidia.com "
     "to nvidia.github.io/dcgm-exporter/helm-charts.",
     ACCENT_AMBER),

    ("Linode Provider Version",
     "CLAUDE.md spec said ~> 2.0. VPC on LKE requires >= 3.3. "
     "Firewall per pool requires >= 3.8. api_version must be v4beta.",
     "Pin to ~> 3.8 minimum. Set api_version = \"v4beta\" on the provider block.",
     ACCENT_BLUE),

    ("OBJ Storage Endpoint Construction",
     "Hand-building endpoint as $region.linodeobjects.com produced wrong URLs. "
     "The bucket resource exposes s3_endpoint as a computed attribute.",
     "Use linode_object_storage_bucket.*.s3_endpoint, never construct the URL manually.",
     ACCENT_BLUE),

    ("Image Registry Hardcoding",
     "Hardcoded Docker Hub username in manifests caused ImagePullBackOff in all environments "
     "except the original author's. Discovered during first real deployment.",
     "Use PLACEHOLDER tokens in YAML, sed-inject at deploy time. "
     "Registry is a required deploy.sh argument.",
     GREEN),
]

for i, (title, problem, fix, color) in enumerate(lessons):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(2.2 + row * 2.6)
    add_card(slide, x, y, Inches(3.95), Inches(2.4))
    add_textbox(slide, x + Inches(0.15), y + Inches(0.05), Inches(3.6), Inches(0.35),
                title, font_size=12, color=color, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.4), Inches(3.6), Inches(0.9),
                problem, font_size=10, color=TEXT_SECONDARY)
    add_textbox(slide, x + Inches(0.15), y + Inches(1.4), Inches(0.5), Inches(0.2),
                "FIX:", font_size=9, color=GREEN, bold=True)
    add_textbox(slide, x + Inches(0.65), y + Inches(1.4), Inches(3.1), Inches(0.9),
                fix, font_size=10, color=TEXT_PRIMARY)

# ----------------------------------------------------------
# SLIDE 13: Security Model
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
section_header(slide, "SECURITY", "Access Control + Secrets Management")

# Left: What's locked down
add_textbox(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(0.3),
            "Current Controls (v1)", font_size=18, color=TEXT_PRIMARY, bold=True)

controls = [
    ("Cloud Firewall", "All external ports (80, 443, 6443, 3000) restricted to operator_ips CIDRs"),
    ("VPC Isolation", "All inter-node traffic on private 10.0.0.0/16 subnet"),
    ("GPU Taint", "NoSchedule prevents non-inference workloads on GPU nodes"),
    ("OBJ Storage ACL", "Bucket is private. Reader key = init containers only. Writer key = operator only"),
    ("Secrets Flow", "TF output \u2192 deploy.sh env var \u2192 kubectl create secret. Never in YAML files"),
    ("Grafana Auth", "random_password resource, anonymous viewer for iframe embedding"),
    ("Kubeconfig", "Sensitive local file, gitignored, chmod 600"),
]

for i, (k, v) in enumerate(controls):
    y = Inches(2.6 + i * 0.55)
    add_textbox(slide, Inches(0.8), y, Inches(2.0), Inches(0.3),
                k, font_size=12, color=ACCENT, bold=True, font_name=FONT_MONO)
    add_textbox(slide, Inches(3.0), y, Inches(3.5), Inches(0.5),
                v, font_size=12, color=TEXT_PRIMARY)

# Right: v2 recommendations
add_card(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.5))
add_textbox(slide, Inches(7.3), Inches(2.3), Inches(4), Inches(0.3),
            "V2 HARDENING PATH", font_size=12, color=ACCENT_AMBER, bold=True)

v2_items = [
    "HashiCorp Vault or External Secrets Operator for centralized secret management",
    "TLS termination on NodeBalancer (cert-manager + Let's Encrypt)",
    "Network Policies for pod-to-pod traffic restriction",
    "Ray token authentication via KubeRay 1.6.0 authOptions",
    "PodSecurityStandards (restricted) on inference namespace",
    "RBAC: least-privilege ServiceAccounts per workload",
    "Audit logging to Akamai ACLP (available in in-bom-2)",
    "VPN-only access to Grafana (remove LoadBalancer, use ClusterIP + VPN)",
]
add_bullet_list(slide, Inches(7.3), Inches(2.7), Inches(5), Inches(3.8),
                v2_items, font_size=12, color=TEXT_SECONDARY, bullet_color=ACCENT_AMBER)

# ----------------------------------------------------------
# SLIDE 14: Performance Characteristics
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
section_header(slide, "PERFORMANCE", "Latency Budget + Resource Utilization")

# Latency breakdown
add_textbox(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(0.3),
            "Inference Latency Budget (p50 estimate)", font_size=18, color=TEXT_PRIMARY, bold=True)

latency_rows = [
    ("Client \u2192 NodeBalancer", "5-15ms", "Depends on client location"),
    ("NodeBalancer \u2192 FastAPI", "<1ms", "Same VPC"),
    ("FastAPI \u2192 Ray Serve", "<1ms", "Cluster-internal service"),
    ("Ray routing", "1-3ms", "Proxy \u2192 actor dispatch"),
    ("CLIP text inference", "10-30ms", "GPU, batch=1, 768-dim output"),
    ("CLIP image inference", "20-50ms", "GPU, preprocessing + forward"),
    ("CLAP audio inference", "30-80ms", "GPU, resampling + forward"),
    ("Return path", "2-5ms", "Serialization + HTTP response"),
]

for i, (stage, time, note) in enumerate(latency_rows):
    y = Inches(2.6 + i * 0.42)
    bg = BG_CARD if i % 2 == 0 else BG_DARK
    add_card(slide, Inches(0.8), y, Inches(5.5), Inches(0.38), bg)
    add_textbox(slide, Inches(0.9), y + Inches(0.03), Inches(2.2), Inches(0.3),
                stage, font_size=11, color=TEXT_PRIMARY)
    add_textbox(slide, Inches(3.2), y + Inches(0.03), Inches(1.0), Inches(0.3),
                time, font_size=11, color=ACCENT, bold=True, font_name=FONT_MONO)
    add_textbox(slide, Inches(4.3), y + Inches(0.03), Inches(1.8), Inches(0.3),
                note, font_size=10, color=TEXT_MUTED)

# Right: Resource allocation
add_textbox(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(0.3),
            "Resource Allocation per Component", font_size=18, color=TEXT_PRIMARY, bold=True)

resources = [
    ("Ray Head", "2-4 CPU, 4-8Gi RAM, 0 GPU", "Scheduler only"),
    ("Ray Worker (x3)", "2-4 CPU, 12-14Gi RAM, 1 GPU", "CLIP or CLAP actor"),
    ("Demo App (x2)", "250m-1 CPU, 256-512Mi RAM", "FastAPI + static files"),
    ("Prometheus", "From kube-prometheus-stack defaults", "~1Gi RAM typical"),
    ("Grafana", "From kube-prometheus-stack defaults", "~256Mi RAM typical"),
    ("DCGM Exporter (x3)", "100-500m CPU, 512Mi-1Gi RAM", "Per GPU node"),
    ("Device Plugin (x3)", "From upstream DaemonSet defaults", "Minimal footprint"),
]

for i, (comp, alloc, note) in enumerate(resources):
    y = Inches(2.6 + i * 0.6)
    add_card(slide, Inches(7.0), y, Inches(5.5), Inches(0.5))
    add_textbox(slide, Inches(7.2), y + Inches(0.02), Inches(2.0), Inches(0.25),
                comp, font_size=11, color=ACCENT_AMBER, bold=True, font_name=FONT_MONO)
    add_textbox(slide, Inches(7.2), y + Inches(0.25), Inches(3.0), Inches(0.25),
                alloc, font_size=10, color=TEXT_PRIMARY, font_name=FONT_MONO)
    add_textbox(slide, Inches(10.3), y + Inches(0.1), Inches(2.0), Inches(0.3),
                note, font_size=10, color=TEXT_MUTED)

# ----------------------------------------------------------
# SLIDE 15: What's Next
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
section_header(slide, "ROADMAP", "What's Next")

tracks = [
    ("NEAR-TERM ENHANCEMENTS", [
        "Zero-shot classification demo: image + category list \u2192 top-10 ranked labels",
        "PersistentVolume model cache to avoid re-pull on pod restart",
        "TLS termination on NodeBalancer via cert-manager",
        "Autoscaling config in serveConfigV2 (min/max replicas based on queue depth)",
        "Helm chart packaging for the full stack (single values.yaml)",
    ], ACCENT),
    ("MULTI-REGION EXPANSION", [
        "Terraform workspaces: one per region, same codebase",
        "Akamai GTM for latency-based DNS routing + automatic failover",
        "Independent Ray clusters per region (no cross-region federation)",
        "Per-region Object Storage with model weight replication",
        "Thanos for federated cross-region Prometheus aggregation",
    ], ACCENT_AMBER),
    ("PRODUCTION HARDENING", [
        "Vault / External Secrets Operator for secrets management",
        "Ray token authentication via KubeRay authOptions",
        "Network Policies + PodSecurityStandards",
        "CI/CD pipeline: build \u2192 push \u2192 deploy with image digest pinning",
        "Load testing with Locust to establish p99 baselines",
    ], ACCENT_BLUE),
]

for i, (title, items, color) in enumerate(tracks):
    x = Inches(0.5 + i * 4.15)
    add_card(slide, x, Inches(2.2), Inches(3.95), Inches(5.0))
    add_textbox(slide, x + Inches(0.2), Inches(2.3), Inches(3.5), Inches(0.3),
                title, font_size=11, color=color, bold=True)
    add_accent_line(slide, x + Inches(0.2), Inches(2.65), Inches(1.5))
    add_bullet_list(slide, x + Inches(0.2), Inches(2.8), Inches(3.5), Inches(4.0),
                    items, font_size=12, color=TEXT_PRIMARY, bullet_color=color)

# ----------------------------------------------------------
# SAVE
# ----------------------------------------------------------
output_path = "/home/morduun/projects/akamai/ray-dist-gpu/CLIP-CLAP-Technical-Deep-Dive.pptx"
prs.save(output_path)
print(f"Deck saved to {output_path}")
